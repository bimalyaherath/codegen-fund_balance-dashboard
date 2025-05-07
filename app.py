import streamlit as st
import pandas as pd
import datetime as dt
import glob
import os
import requests
import json
import tempfile
import io

# Try imports for export functionality
try:
    from fpdf import FPDF
    from pptx import Presentation
    from pptx.util import Inches
    EXPORT_ENABLED = True
except ImportError:
    EXPORT_ENABLED = False

# --- Data Upload & Versioning Setup ---
UPLOAD_DIR = os.path.join(tempfile.gettempdir(), 'weekly_uploads')
HISTORY_FILE = os.path.join(UPLOAD_DIR, 'upload_history.json')
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Load or initialize upload history
def load_history():
    if os.path.exists(HISTORY_FILE):
        try:
            return json.load(open(HISTORY_FILE))
        except:
            return []
    return []

history = load_history()

# File uploader for versioned weekly files
uploaded_files = st.sidebar.file_uploader(
    'Upload Weekly Excel Files (versioned)',
    type=['xlsx'], accept_multiple_files=True
)
if uploaded_files:
    timestamp = dt.datetime.now().isoformat()
    saved = []
    for f in uploaded_files:
        filename = f"{timestamp}_{f.name}"
        path = os.path.join(UPLOAD_DIR, filename)
        with open(path, 'wb') as out:
            out.write(f.getbuffer())
        saved.append(filename)
    history.append({'timestamp': timestamp, 'files': saved})
    json.dump(history, open(HISTORY_FILE, 'w'))
    st.sidebar.success(f"Uploaded {len(saved)} files at {timestamp}")

# Version selection
version_labels = ['Latest'] + [h['timestamp'] for h in history]
version_choice = st.sidebar.selectbox('Select Data Version', version_labels)
if version_choice == 'Latest':
    version_files = os.listdir(UPLOAD_DIR)
else:
    version_files = next(h['files'] for h in history if h['timestamp'] == version_choice)

# Load and parse data
@st.cache_data
def load_fund_data(files):
    static_paths = glob.glob('data/Fund_Balance_*.xlsx') \
                 + glob.glob('/mnt/data/Fund_Balance_*.xlsx') \
                 + glob.glob('Fund_Balance_*.xlsx')
    upload_paths = [os.path.join(UPLOAD_DIR, f) for f in files]
    all_paths = static_paths + upload_paths
    frames = []
    for file in all_paths:
        week_label = os.path.basename(file).replace('Fund_Balance_','').replace('.xlsx','')
        try:
            df = pd.read_excel(file)
        except:
            continue
        mask = df['Details'].str.contains('Bank & Cash Balances|Cash Ins|Cash Outs', na=False)
        idx = df[mask].index.tolist()
        parts = []
        for i, start in enumerate(idx):
            end = idx[i+1] if i+1 < len(idx) else len(df)
            section = df.loc[start, 'Details']
            part = df.iloc[start+1:end].dropna(how='all').copy()
            part['Section'] = section
            parts.append(part)
        week_df = pd.concat(parts, ignore_index=True)
        week_df['Week'] = week_label
        frames.append(week_df)
    return pd.concat(frames, ignore_index=True) if frames else pd.DataFrame()

fund_data = load_fund_data(version_files)
if fund_data.empty:
    st.error('No fund data found for selected version.')
    st.stop()

# Helper: parse week start date
def parse_week(w):
    try:
        start = w.split('_to_')[0]
        return dt.datetime.strptime(start, '%B_%d').replace(year=dt.date.today().year)
    except:
        return dt.datetime.min

weeks = sorted(fund_data['Week'].unique(), key=parse_week)

# --- Sidebar Filters ---
st.sidebar.title('Dashboard Settings')
selected_week = st.sidebar.selectbox('Select Week', weeks)

currencies = ['LKR','USD','GBP','AUD','DKK','EUR','MXN','INR','AED']
selected_currencies = st.sidebar.multiselect('Select Currencies', currencies, default=[currencies[0]])
if not selected_currencies:
    st.sidebar.error('Select at least one currency')
    st.stop()

# Live currency converter
@st.cache_data(ttl=3600)
def get_rate(base, symbol):
    resp = requests.get('https://api.exchangerate.host/latest', params={'base': base, 'symbols': symbol}, timeout=5)
    return resp.json().get('rates', {}).get(symbol)

st.sidebar.header('Currency Converter')
amt = st.sidebar.number_input('Amount', value=1.0)
frm = st.sidebar.selectbox('From Currency', currencies)
to_cur = st.sidebar.selectbox('To Currency', currencies, index=1)
r = get_rate(frm, to_cur)
if r:
    st.sidebar.write(f'1 {frm} = {r:.4f} {to_cur}')
    st.sidebar.write(f'{amt} {frm} = {amt*r:.2f} {to_cur}')
else:
    st.sidebar.error('Failed to fetch live exchange rate')

# --- Main Dashboard ---
try:
    sd = parse_week(selected_week).date()
    ed = sd + dt.timedelta(days=6)
    subtitle = f"{selected_week.replace('_',' ')} ({sd} to {ed})"
except:
    subtitle = selected_week.replace('_',' ')
st.title('📊 Weekly Fund Dashboard')
st.subheader(subtitle)

# Filter week data
df_week = fund_data[fund_data['Week'] == selected_week]

# 1. Weekly Summary
st.header('🔖 Weekly Summary')
sum_df = df_week.groupby('Section')[selected_currencies].sum().reset_index()
sum_df.columns = ['Category'] + [f'Total {c}' for c in selected_currencies]
st.table(sum_df)
st.download_button('Download Weekly Summary', sum_df.to_csv(index=False), file_name=f'{selected_week}_Weekly_Summary.csv', mime='text/csv')

# 2. Cash Breakdown
st.header('📂 Cash Ins & Outs Breakdown')
with st.expander('Cash Ins'):
    ins_tbl = df_week[df_week['Section']=='Cash Ins'][['Details'] + selected_currencies].rename(columns={'Details':'Category'})
    st.table(ins_tbl)
with st.expander('Cash Outs'):
    outs_tbl = df_week[df_week['Section']=='Cash Outs'][['Details'] + selected_currencies].rename(columns={'Details':'Category'})
    st.table(outs_tbl)

# 3. Weekly Comparison
st.header('📈 Weekly Comparison')
comp_weeks = st.multiselect('Select Weeks to Compare', weeks, default=weeks[-2:])
if comp_weeks:
    comp_df = fund_data[fund_data['Week'].isin(comp_weeks)]
    comp_summary = comp_df.groupby('Week')[selected_currencies].sum()
    st.line_chart(comp_summary)
else:
    st.info('Select at least one week to compare.')

# 4. Additional Charts
st.header('📊 Additional Charts')
st.subheader('Cash Ins vs Cash Outs')
st.bar_chart(df_week.groupby('Section')[selected_currencies].sum().loc[['Cash Ins','Cash Outs']])

# 5. Net Cash-Flow & Balances
st.header('💰 Net Cash-Flow & Balances')
for c in selected_currencies:
    open_vals = fund_data[fund_data['Section']=='Bank & Cash Balances'].groupby('Week')[c].sum()
    ins_vals = fund_data[fund_data['Section']=='Cash Ins'].groupby('Week')[c].sum()
    outs_vals = fund_data[fund_data['Section']=='Cash Outs'].groupby('Week')[c].sum()
    net_vals = ins_vals.sub(outs_vals, fill_value=0)
    close_vals = open_vals.add(net_vals, fill_value=0)
    df_bal = pd.DataFrame({'Opening': open_vals, 'Net Change': net_vals, 'Closing': close_vals}).reindex(weeks)
    st.subheader(f'{c} Balances Over Time')
    st.line_chart(df_bal)

# 6. Alerts & Thresholds
st.header('🚨 Alerts & Thresholds')
thresholds = {c: st.number_input(f'Threshold for net cash change in {c}', value=0.0, step=1.0, key=f'th_{c}') for c in selected_currencies}
ins_w = fund_data[(fund_data['Week']==selected_week)&(fund_data['Section']=='Cash Ins')].groupby('Week')[selected_currencies].sum().reindex([selected_week], fill_value=0)
outs_w = fund_data[(fund_data['Week']==selected_week)&(fund_data['Section']=='Cash Outs')].groupby('Week')[selected_currencies].sum().reindex([selected_week], fill_value=0)
net_w = ins_w.sub(outs_w, fill_value=0)
for c in selected_currencies:
    val = float(net_w[c]); th = thresholds[c]
    if val < th:
        st.error(f"🚨 Alert! Net cash for {selected_week} in {c} is {val:.2f}, below your threshold of {th:.2f}. 📉 Please review.")
    else:
        st.success(f"✅ Net cash for {selected_week} in {c} is {val:.2f}, above threshold {th:.2f}. 🎉 Good job!")

# 7. Date-Range & Rolling Periods
st.header('📅 Date-Range & Rolling Periods')
mode = st.radio('Select Mode:', ['Custom Date Range','Rolling Window'])
mdates = [parse_week(w).date() for w in weeks]
min_date, max_date = min(mdates), max(mdates)
if mode == 'Custom Date Range':
    dr = st.date_input('Date range:', [min_date, max_date], min_value=min_date, max_value=max_date)
    if len(dr)==2:
        start_d, end_d = dr
        selw = [w for w in weeks if start_d <= parse_week(w).date() <= end_d]
        if selw:
            df2 = fund_data[fund_data['Week'].isin(selw)]
            sum2 = df2.groupby('Section')[selected_currencies].sum().reset_index()
            sum2.columns = ['Category'] + [f'Total {c}' for c in selected_currencies]
            st.write(sum2)
        else:
            st.info('No data in this date range.')
elif mode == 'Rolling Window':
    rw = st.selectbox('Window Type:', ['Last N Weeks','Month-to-Date','Quarter-to-Date'])
    if rw == 'Last N Weeks':
        n = st.number_input('Number of weeks:', min_value=1, max_value=len(weeks), value=4)
        selw = weeks[-int(n):]
    elif rw == 'Month-to-Date':
        today = dt.date.today(); start_m = today.replace(day=1)
        selw = [w for w in weeks if start_m <= parse_week(w).date() <= today]
    else:
        today = dt.date.today(); q=(today.month-1)//3; start_q = dt.date(today.year,q*3+1,1)
        selw = [w for w in weeks if start_q <= parse_week(w).date() <= today]
    df2 = fund_data[fund_data['Week'].isin(selw)]
    sum2 = df2.groupby('Section')[selected_currencies].sum().reset_index()
    sum2.columns = ['Category'] + [f'Total {c}' for c in selected_currencies]
    st.write(sum2)

# 8. Forecasting & Trends
# Attempt to import statsmodels for forecasting
try:
    import statsmodels.api as sm
    SM_ENABLED = True
except ImportError:
    SM_ENABLED = False

st.header('🔮 Forecasting & Trends')
if SM_ENABLED:
    # A. Moving Averages on Weekly Comparison
    st.subheader('Moving Averages on Weekly Comparison')
    if 'comp_weeks' in locals() and comp_weeks:
        comp_df = fund_data[fund_data['Week'].isin(comp_weeks)].groupby('Week')[selected_currencies].sum()
        for c in selected_currencies:
            series = comp_df[c]
            ma = series.rolling(window=3, min_periods=1).mean()
            df_plot = pd.DataFrame({'Actual': series, 'Moving Average (3wk)': ma})
            st.line_chart(df_plot)
    else:
        st.info('Select weeks above to show moving averages.')
    # B. ARIMA Forecast for Closing Balance
    st.subheader('ARIMA Forecast for Closing Balance')
    for c in selected_currencies:
        open_vals = fund_data[fund_data['Section']=='Bank & Cash Balances'].groupby('Week')[c].sum()
        ins_vals = fund_data[fund_data['Section']=='Cash Ins'].groupby('Week')[c].sum()
        outs_vals = fund_data[fund_data['Section']=='Cash Outs'].groupby('Week')[c].sum()
        net_vals = ins_vals.sub(outs_vals, fill_value=0)
        close_vals = open_vals.add(net_vals, fill_value=0)
        try:
            model = sm.tsa.ARIMA(close_vals, order=(1,1,0))
            res = model.fit()
            fcast = res.forecast(steps=1)
            st.write(f'Forecast next week closing balance in {c}: {fcast.iloc[0]:.2f} {c}')
        except Exception as e:
            st.warning(f'Could not forecast for {c}: {e}')
else:
    st.info('Install `statsmodels` to enable forecasting & trends')

# 9. Export & Sharing
st.header('📤 Export & Sharing')
if EXPORT_ENABLED:
    if st.button('Export Weekly Summary as PDF'):
        pdf = FPDF(); pdf.add_page(); pdf.set_font('Arial','',10)
        for _, row in sum_df.iterrows():
            pdf.cell(0,8,' | '.join(str(v) for v in row.values), ln=1)
        data = pdf.output(dest='S').encode('latin1')
        st.download_button('Download PDF', data, 'summary.pdf', 'application/pdf')
    if st.button('Export Weekly Summary as PPT'):
        prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[5])
        txt = slide.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(5)).text_frame
        for _, row in sum_df.iterrows(): txt.add_paragraph(' | '.join(str(v) for v in row.values))
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        st.download_button('Download PPT', buf.read(), 'summary.pptx', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
else:
    st.info('Install fpdf and python-pptx to enable exports')

# 10. Full Dataset
st.header('📁 Full Dataset')
# Note about filtering and download
st.info('You can filter the dataset using the search box below and download the resulting subset.')
with st.expander('View & Filter Full Dataset'):
    # Text filter on Details column
    filter_text = st.text_input('Filter by Details (case-insensitive)', '')
    if filter_text:
        filtered_df = fund_data[fund_data['Details'].str.contains(filter_text, case=False, na=False)]
    else:
        filtered_df = fund_data
    st.dataframe(filtered_df)
    st.download_button(
        label='Download Filtered Dataset',
        data=filtered_df.to_csv(index=False),
        file_name='Filtered_Weekly_Fund_Data.csv',
        mime='text/csv'
    )

# Footer
st.write('---')
st.caption('Created with ❤️ using Streamlit & GitHub')
